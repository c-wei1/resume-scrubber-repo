"""Generate a single-slide summary of PII Detector rules (white + red-accent theme)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (same as main deck) ───────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF3, 0xF4, 0xF6)
MID_GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY   = RGBColor(0x4B, 0x55, 0x63)
NEAR_BLACK  = RGBColor(0x1F, 0x23, 0x37)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_LIGHT   = RGBColor(0xFE, 0xE2, 0xE2)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ── Shape helpers (same as main deck) ─────────────────────────
def _bg(slide, c=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def _rect(sl, l, t, w, h, fill=None, line_c=None, line_w=Pt(1)):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line_c: s.line.color.rgb = line_c; s.line.width = line_w
    else: s.line.fill.background()
    s.shadow.inherit = False; return s

def _rrect(sl, l, t, w, h, fill=LIGHT_GRAY, line_c=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if line_c: s.line.color.rgb = line_c; s.line.width = Pt(1.5)
    s.shadow.inherit = False; s.adjustments[0] = 0.03; return s

def _oval(sl, l, t, w, h, fill=RED):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False; return s

def _tb(sl, l, t, w, h):
    return sl.shapes.add_textbox(l, t, w, h)

def _set(tf, txt, sz=18, bold=False, col=NEAR_BLACK, align=PP_ALIGN.LEFT, after=Pt(4)):
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = col; p.alignment = align; p.space_after = after
    p.font.name = "Calibri"; return p

def _p(tf, txt, sz=15, bold=False, col=NEAR_BLACK, align=PP_ALIGN.LEFT,
       level=0, after=Pt(4), before=Pt(0)):
    p = tf.add_paragraph(); p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = col; p.alignment = align; p.level = level
    p.space_after = after; p.space_before = before; p.font.name = "Calibri"; return p

def _line(sl, l, t, w):
    return _rect(sl, l, t, w, Pt(2.5), fill=RED)

def _sidebar(sl, w=Inches(0.12)):
    _rect(sl, Inches(0), Inches(0), w, SH, fill=RED)


# ═══════════════════════════════════════════════════════════════
#  PII DETECTOR RULES — SINGLE SLIDE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
_sidebar(s)

# Title
tb = _tb(s, Inches(0.8), Inches(0.35), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "PII Detection — Regex Rules", sz=36, bold=True)
_line(s, Inches(0.8), Inches(0.95), Inches(2.8))

# ── 5 cards in a row ──────────────────────────────────────────
card_data = [
    ("Phone Number", RED, [
        "Lookbehind: (?<![\\w@.])",
        "Intl prefix: (?:(?:\\+|00)\\s?\\d{1,3}[\\s.\\-]?)?",
        "Area code: (?:\\(\\s?\\d{1,4}\\s?\\)[\\s.\\-]?)?",
        "Core digits: \\d{1,4}(?:[\\s.\\-]?\\d{1,4}){1,5}",
        "Extension: (?:\\s?(?:ext|x)\\.?\\s?\\d{1,5})?",
        "Lookahead: (?![\\w])",
        "",
        "Digit gate: 7 ≤ len(digits) ≤ 15",
    ]),
    ("Email Address", RED, [
        "Local part: [a-zA-Z0-9._%+\\-]+",
        "At sign: @",
        "Domain: [a-zA-Z0-9.\\-]+",
        "Dot: \\.",
        "TLD: [a-zA-Z]{2,}",
        "",
        "",
        "(no validation gate needed)",
    ]),
    ("Web Links / URLs", RED, [
        "Protocol: (?:https?://|www\\.)",
        "Bare domain: [a-zA-Z0-9\\-]+\\.",
        "Known TLDs: com|org|net|io|dev|me|co|info|biz",
        "Boundary: (?:/|\\b)",
        "Path chars: [^\\s<>()\\[\\]\"']+",
        "",
        "Hyperlinks: <w:hyperlink> unwrapped",
        "Rels: image+hyperlink types stripped",
    ]),
    ("Images", RED, [
        "Containers: //w:drawing | //w:pict",
        "Text check: .//*[local-name()='t']",
        "Textbox check: .//*[local-name()='txbx']",
        "",
        "Has text: strip pic/blip/blipFill/imagedata",
        "No text: remove entire element",
        "",
        "Rel type: …/relationships/image removed",
    ]),
    ("Addresses", RED, [
        "Postal codes: \\b\\d{5}(?:-\\d{4})?\\b  (+3)",
        "Street type: suffix|prefix|fused  (+3)",
        "PO Box: p\\.?o\\.?\\s*box  (+3)",
        "House num: \\d{1,5}[A-Za-z]?  (+2)",
        "Region: AL|AK|…|DC|NSW|…  (+2)",
        "Weak postal: \\d{4}\\s+[A-Z]\\w+  (+2)",
        "Unit/suite: suite|apt|unit\\s*\\w+  (+1)",
        "Direction: \\b[NSEW]{1,2}\\.?\\s  (+1)",
        "Country: USA|United States|…  (+1)",
    ]),
]

card_w = Inches(2.35)
card_h = Inches(4.6)
gap = Inches(0.15)
start_x = Inches(0.65)
card_y = Inches(1.25)

for i, (title, accent, items) in enumerate(card_data):
    x = start_x + i * (card_w + gap)

    # Card background
    _rrect(s, x, card_y, card_w, card_h, LIGHT_GRAY)
    # Accent header bar
    _rrect(s, x, card_y, card_w, Inches(0.42), accent)
    # Title text
    tb = _tb(s, x, card_y + Inches(0.06), card_w, Inches(0.34))
    _set(tb.text_frame, title, sz=12, bold=True, col=WHITE,
         align=PP_ALIGN.CENTER, after=Pt(0))
    # Bullet items (monospace-style for regex)
    for j, item in enumerate(items):
        tb = _tb(s, x + Inches(0.1), card_y + Inches(0.48 + j * 0.44),
                 card_w - Inches(0.2), Inches(0.4))
        tf = tb.text_frame; tf.word_wrap = True
        p = _set(tf, item, sz=9, col=DARK_GRAY if item else WHITE, after=Pt(0))
        if item:
            p.font.name = "Consolas"

# Footer: threshold note for addresses
tb = _tb(s, Inches(0.8), Inches(6.2), Inches(11.0), Inches(0.6))
tf = tb.text_frame; tf.word_wrap = True
_set(tf, "Addresses: score ≥ 5 → redact entire line  •  Neighbour context boost  •  "
     "PII (phone/email/URL) stripped before address scoring to avoid false positives",
     sz=10, col=MID_GRAY, after=Pt(0))

# ── Save ──────────────────────────────────────────────────────
prs.save("pii_detection_rules.pptx")
print("Saved → pii_detection_rules.pptx")
