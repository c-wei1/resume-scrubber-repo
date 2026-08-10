"""Generate a single-slide 'About Me' intro in the same white + red-accent theme."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (same as main deck) ───────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF3, 0xF4, 0xF6)
MID_GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY   = RGBColor(0x4B, 0x55, 0x63)
NEAR_BLACK  = RGBColor(0x1F, 0x23, 0x37)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_LIGHT   = RGBColor(0xFE, 0xE2, 0xE2)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ── Shape helpers (same as main deck) ─────────────────────────
def _bg(slide, c=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def _rect(sl, l, t, w, h, fill=None, line_c=None, line_w=Pt(1)):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line_c: s.line.color.rgb = line_c; s.line.width = line_w
    else: s.line.fill.background()
    s.shadow.inherit = False; return s

def _rrect(sl, l, t, w, h, fill=LIGHT_GRAY, line_c=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if line_c: s.line.color.rgb = line_c; s.line.width = Pt(1.5)
    s.shadow.inherit = False; s.adjustments[0] = 0.03; return s

def _oval(sl, l, t, w, h, fill=RED):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False; return s

def _tb(sl, l, t, w, h):
    return sl.shapes.add_textbox(l, t, w, h)

def _set(tf, txt, sz=18, bold=False, col=NEAR_BLACK, align=PP_ALIGN.LEFT, after=Pt(4)):
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = col; p.alignment = align; p.space_after = after
    p.font.name = "Calibri"; return p

def _p(tf, txt, sz=15, bold=False, col=NEAR_BLACK, align=PP_ALIGN.LEFT,
       level=0, after=Pt(4), before=Pt(0)):
    p = tf.add_paragraph(); p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = col; p.alignment = align; p.level = level
    p.space_after = after; p.space_before = before; p.font.name = "Calibri"; return p

def _line(sl, l, t, w):
    return _rect(sl, l, t, w, Pt(2.5), fill=RED)

def _sidebar(sl, w=Inches(0.12)):
    _rect(sl, Inches(0), Inches(0), w, SH, fill=RED)


# ═══════════════════════════════════════════════════════════════
#  ABOUT ME — INTRODUCTION SLIDE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
_sidebar(s)

# ── Title ──────────────────────────────────────────────────────
tb = _tb(s, Inches(0.8), Inches(0.35), Inches(8.0), Inches(0.7))
_set(tb.text_frame, "About Me", sz=36, bold=True)
_line(s, Inches(0.8), Inches(0.95), Inches(2.0))

# ── Decorative accent circles (top-right, matching title slide) ──
_oval(s, Inches(10.8), Inches(0.3), Inches(2.0), Inches(2.0), RED_LIGHT)
_oval(s, Inches(11.3), Inches(0.8), Inches(1.0), Inches(1.0), RED)

# ── Name / headline area ──────────────────────────────────────
tb = _tb(s, Inches(0.8), Inches(1.2), Inches(6.0), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
_set(tf, "Cynthia Wei", sz=28, bold=True, col=RED, after=Pt(2))
_p(tf, "4th-Year Undergraduate  ·  Computer Science  ·  UCLA", sz=14, col=DARK_GRAY, after=Pt(0))

# ── Three info cards ───────────────────────────────────────────
card_w = Inches(3.7)
card_h = Inches(3.8)
gap    = Inches(0.2)
cards_total_w = 3 * card_w + 2 * gap
start_x = (SW - cards_total_w) // 2
card_y = Inches(2.5)

# --- Card 1: Background ---
x1 = start_x
_rrect(s, x1, card_y, card_w, card_h, LIGHT_GRAY)
_rect(s, x1, card_y, card_w, Pt(4), fill=RED)

tb = _tb(s, x1 + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.4))
_set(tb.text_frame, "Background", sz=16, bold=True, col=RED, after=Pt(6))

items_bg = [
    ("Education", "B.S. Computer Science, UCLA\n4th year (expected 2027)"),
    ("Hometown", "Bay Area, California"),
]
for j, (lbl, val) in enumerate(items_bg):
    iy = card_y + Inches(0.7 + j * 1.2)
    _oval(s, x1 + Inches(0.25), iy + Inches(0.05), Inches(0.13), Inches(0.13), RED)
    tb = _tb(s, x1 + Inches(0.5), iy - Inches(0.05), card_w - Inches(0.7), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, lbl, sz=13, bold=True, after=Pt(2))
    _p(tf, val, sz=12, col=DARK_GRAY, after=Pt(0))

# --- Card 2: Hobbies & Interests ---
x2 = start_x + card_w + gap
_rrect(s, x2, card_y, card_w, card_h, LIGHT_GRAY)
_rect(s, x2, card_y, card_w, Pt(4), fill=RED)

tb = _tb(s, x2 + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.4))
_set(tb.text_frame, "Hobbies & Interests", sz=16, bold=True, col=RED, after=Pt(6))

hobbies = [
    ("Hiking", "Exploring trails and nature"),
    ("Reading", "Always looking for the next great book"),
    ("Baking", "Experimenting with new recipes"),
]
for j, (lbl, val) in enumerate(hobbies):
    iy = card_y + Inches(0.7 + j * 0.95)
    _oval(s, x2 + Inches(0.25), iy + Inches(0.05), Inches(0.13), Inches(0.13), RED)
    tb = _tb(s, x2 + Inches(0.5), iy - Inches(0.05), card_w - Inches(0.7), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, lbl, sz=13, bold=True, after=Pt(2))
    _p(tf, val, sz=12, col=DARK_GRAY, after=Pt(0))

# --- Card 3: Career & Experience ---
x3 = start_x + 2 * (card_w + gap)
_rrect(s, x3, card_y, card_w, card_h, LIGHT_GRAY)
_rect(s, x3, card_y, card_w, Pt(4), fill=RED)

tb = _tb(s, x3 + Inches(0.2), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.4))
_set(tb.text_frame, "Career & Experience", sz=16, bold=True, col=RED, after=Pt(6))

career = [
    ("Data Science Union", "Member of UCLA's premier data\nscience organization"),
    ("Alignment Research", "Researching AI alignment —\nensuring AI systems act in\naccordance with human values"),
]
for j, (lbl, val) in enumerate(career):
    iy = card_y + Inches(0.7 + j * 1.35)
    _oval(s, x3 + Inches(0.25), iy + Inches(0.05), Inches(0.13), Inches(0.13), RED)
    tb = _tb(s, x3 + Inches(0.5), iy - Inches(0.05), card_w - Inches(0.7), Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, lbl, sz=13, bold=True, after=Pt(2))
    _p(tf, val, sz=12, col=DARK_GRAY, after=Pt(0))


# ── Save ───────────────────────────────────────────────────────
out = "Introduction_Slide.pptx"
prs.save(out)
print(f"Saved → {out}")
