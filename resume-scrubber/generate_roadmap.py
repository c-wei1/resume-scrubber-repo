"""Generate a compact 2-slide 'Next Steps' deck — timeline + 6-month metrics."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ───────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF3, 0xF4, 0xF6)
MID_GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY   = RGBColor(0x4B, 0x55, 0x63)
NEAR_BLACK  = RGBColor(0x1F, 0x23, 0x37)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_LIGHT   = RGBColor(0xFE, 0xE2, 0xE2)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ── Shape helpers ─────────────────────────────────────────────
def _bg(slide, c=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def _rect(sl, l, t, w, h, fill=None, line_c=None, line_w=Pt(1)):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line_c: s.line.color.rgb = line_c; s.line.width = line_w
    else: s.line.fill.background()
    s.shadow.inherit = False; return s

def _rrect(sl, l, t, w, h, fill=LIGHT_GRAY, line_c=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if line_c: s.line.color.rgb = line_c; s.line.width = Pt(1.5)
    s.shadow.inherit = False; s.adjustments[0] = 0.03; return s

def _oval(sl, l, t, w, h, fill=RED):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False; return s

def _tb(sl, l, t, w, h):
    return sl.shapes.add_textbox(l, t, w, h)

def _set(tf, txt, sz=18, bold=False, col=NEAR_BLACK, align=PP_ALIGN.LEFT, after=Pt(4)):
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = col; p.alignment = align; p.space_after = after
    p.font.name = "Calibri"; return p

def _p(tf, txt, sz=15, bold=False, col=NEAR_BLACK, align=PP_ALIGN.LEFT,
       level=0, after=Pt(4), before=Pt(0)):
    p = tf.add_paragraph(); p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = col; p.alignment = align; p.level = level
    p.space_after = after; p.space_before = before; p.font.name = "Calibri"; return p

def _chevron(sl, l, t, w, h, fill=RED):
    s = sl.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False; return s

def _line(sl, l, t, w):
    return _rect(sl, l, t, w, Pt(2.5), fill=RED)

def _sn(sl, n):
    tb = _tb(sl, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4))
    _set(tb.text_frame, str(n), sz=10, col=MID_GRAY, align=PP_ALIGN.RIGHT)

def _sidebar(sl, w=Inches(0.12)):
    _rect(sl, Inches(0), Inches(0), w, SH, fill=RED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — 6-MONTH ROADMAP TIMELINE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Next Steps — 6-Month Roadmap", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(3.6))

# Timeline bar
_rect(s, Inches(0.8), Inches(2.05), Inches(11.7), Pt(4), fill=MID_GRAY)

phases = [
    ("Months 1–2",  "Strengthen the Core",                      LIGHT_GRAY, NEAR_BLACK),
    ("Months 3–4",  "Smarter Detection &\nCross-Team Alignment", LIGHT_GRAY, NEAR_BLACK),
    ("Months 5–6",  "Integration & Scale",                      RED_LIGHT,  RED),
]

phase_items = [
    [
        "Expand automated test suite & CI/CD\npipeline with integration tests",
        "Batch processing — upload and process\n50+ CVs in a single run",
        "PDF & .doc ingestion via pre-conversion\nbefore the main pipeline",
        "Extend address scorer to CJK, Arabic,\nand Cyrillic address formats",
    ],
    [
        "Fine-tune NER model on a larger, more\ndiverse resume corpus",
        "Collaborate with Kite & PDM training\nteams to harmonize templates & workflows",
        "Post-redaction confidence scoring and\naudit report for reviewer attention",
        "Extend PII scrubbing and template\nconformance to Job Descriptions",
    ],
    [
        "Direct GVault API integration — replace\nWalkMe with end-to-end automation",
        "Real-time progress streaming for batch\noperations via WebSocket / SSE",
        "Monitoring dashboard — processing volume,\nPII detection rates, throughput",
        "Production hardening — structured logging,\nalerting, SLA tracking",
    ],
]

pcw = Inches(3.6); pgap = Inches(0.45)
ptw = len(phases) * pcw + (len(phases) - 1) * pgap
psx = (SW - ptw) // 2

for i, (period, title, bg_c, accent) in enumerate(phases):
    x = psx + i * (pcw + pgap)
    is_last = (i == len(phases) - 1)

    # Timeline dot
    dsz = Inches(0.26)
    dot = _oval(s, x + pcw / 2 - dsz / 2, Inches(1.93), dsz, dsz,
                RED if is_last else WHITE)
    if not is_last:
        dot.line.color.rgb = RED; dot.line.width = Pt(2)

    # Period label
    tb = _tb(s, x, Inches(1.5), pcw, Inches(0.4))
    _set(tb.text_frame, period, sz=12, bold=True, col=RED, align=PP_ALIGN.CENTER, after=Pt(0))

    # Phase card
    card_y = Inches(2.5)
    _rrect(s, x, card_y, pcw, Inches(4.6), bg_c, line_c=RED if is_last else None)

    # Card title
    tb = _tb(s, x + Inches(0.12), card_y + Inches(0.12), pcw - Inches(0.24), Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, title, sz=15, bold=True, col=accent, align=PP_ALIGN.CENTER, after=Pt(0))

    # Items
    for j, item in enumerate(phase_items[i]):
        iy = card_y + Inches(0.85 + j * 0.9)
        _oval(s, x + Inches(0.18), iy + Inches(0.08), Inches(0.12), Inches(0.12), accent)
        tb = _tb(s, x + Inches(0.4), iy - Inches(0.02), pcw - Inches(0.55), Inches(0.85))
        tf = tb.text_frame; tf.word_wrap = True
        _set(tf, item, sz=11, col=DARK_GRAY, after=Pt(0))

_sn(s, 1)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — 6-MONTH TARGET METRICS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "6-Month Target Metrics", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(2.8))

# 2x4 grid of metric cards
metrics = [
    ("50+",      "CVs per batch",          "Process an entire hiring cycle in\na single upload — no more one-by-one"),
    ("90%+",     "test coverage",           "Comprehensive automated test suite\nacross PII, section detection, and API"),
    ("3",        "input formats",           ".docx + PDF + .doc support —\naccept any common resume format"),
    ("2×",       "document types",          "Extend pipeline to Job Descriptions\nin addition to CVs"),
    ("3 teams",  "aligned",                 "Kite, PDM, and QCE on unified\ntemplates and review workflows"),
    ("↑ NER",    "model accuracy",          "Fine-tuned on larger, diverse corpus\nwith measurable precision/recall gains"),
    ("0",        "manual upload steps",     "Direct GVault API integration\nreplaces WalkMe workflow entirely"),
    ("Live",     "monitoring dashboard",    "Real-time processing volume, PII\ndetection rates, and system health"),
]

mcw = Inches(2.75); mh = Inches(2.35)
mg = Inches(0.2)
cols = 4; rows = 2
grid_w = cols * mcw + (cols - 1) * mg
grid_h = rows * mh + (rows - 1) * mg
sx = (SW - grid_w) // 2
sy = Inches(1.5)

for i, (big, lbl, desc) in enumerate(metrics):
    col_idx = i % cols
    row_idx = i // cols
    x = sx + col_idx * (mcw + mg)
    y = sy + row_idx * (mh + mg)

    is_accent = (i in (0, 3, 6))  # highlight key items
    bg_c = RED_LIGHT if is_accent else LIGHT_GRAY
    accent = RED if is_accent else NEAR_BLACK

    _rrect(s, x, y, mcw, mh, bg_c, line_c=RED if is_accent else None)
    _rect(s, x, y, Pt(4), mh, fill=accent)

    # Big number
    tb = _tb(s, x + Inches(0.25), y + Inches(0.15), mcw - Inches(0.4), Inches(0.65))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, big, sz=30, bold=True, col=accent, after=Pt(0))
    _p(tf, lbl, sz=11, col=MID_GRAY, after=Pt(0))

    # Description
    tb = _tb(s, x + Inches(0.25), y + Inches(1.05), mcw - Inches(0.4), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, desc, sz=11, col=DARK_GRAY, after=Pt(0))

# Bottom tagline
tb = _tb(s, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.4))
tf = tb.text_frame; tf.word_wrap = True
_set(tf, "Each phase builds on the last — from a CV scrubber to a full document compliance platform.",
     sz=13, bold=True, col=MID_GRAY, align=PP_ALIGN.CENTER, after=Pt(0))
_sn(s, 2)


# ── Save ───────────────────────────────────────────────────────
out = "Next_Steps_Roadmap.pptx"
prs.save(out)
print(f"Saved → {out}")

