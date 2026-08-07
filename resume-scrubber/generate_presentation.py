"""Generate a minimal white + red-accent presentation for Resume Scrubber."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Palette ────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF3, 0xF4, 0xF6)
MID_GRAY    = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY   = RGBColor(0x4B, 0x55, 0x63)
NEAR_BLACK  = RGBColor(0x1F, 0x23, 0x37)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_LIGHT   = RGBColor(0xFE, 0xE2, 0xE2)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_DARK  = RGBColor(0x16, 0xA3, 0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


# ── Shape helpers ──────────────────────────────────────────────
def _bg(slide, c=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def _rect(sl, l, t, w, h, fill=None, line_c=None, line_w=Pt(1)):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line_c: s.line.color.rgb = line_c; s.line.width = line_w
    else: s.line.fill.background()
    s.shadow.inherit = False; return s

def _rrect(sl, l, t, w, h, fill=LIGHT_GRAY, line_c=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if line_c: s.line.color.rgb = line_c; s.line.width = Pt(1.5)
    s.shadow.inherit = False; s.adjustments[0] = 0.03; return s

def _oval(sl, l, t, w, h, fill=RED):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False; return s

def _tb(sl, l, t, w, h):
    return sl.shapes.add_textbox(l, t, w, h)

def _set(tf, txt, sz=18, bold=False, col=NEAR_BLACK, align=PP_ALIGN.LEFT, after=Pt(4)):
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = col; p.alignment = align; p.space_after = after
    p.font.name = "Calibri"; return p

def _p(tf, txt, sz=15, bold=False, col=NEAR_BLACK, align=PP_ALIGN.LEFT,
       level=0, after=Pt(4), before=Pt(0)):
    p = tf.add_paragraph(); p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = col; p.alignment = align; p.level = level
    p.space_after = after; p.space_before = before; p.font.name = "Calibri"; return p

def _chevron(sl, l, t, w, h, fill=RED):
    s = sl.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False; return s

def _arrow(sl, l, t, w, h, fill=MID_GRAY):
    s = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.shadow.inherit = False; return s

def _line(sl, l, t, w):
    return _rect(sl, l, t, w, Pt(2.5), fill=RED)

def _sn(sl, n):
    tb = _tb(sl, Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4))
    _set(tb.text_frame, str(n), sz=10, col=MID_GRAY, align=PP_ALIGN.RIGHT)

def _sidebar(sl, w=Inches(0.12)):
    _rect(sl, Inches(0), Inches(0), w, SH, fill=RED)


# ═══════════════════════════════════════════════════════════════
# 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
_sidebar(s, Inches(0.18))
_oval(s, Inches(10.5), Inches(0.5), Inches(2.4), Inches(2.4), RED_LIGHT)
_oval(s, Inches(11.1), Inches(1.1), Inches(1.2), Inches(1.2), RED)

tb = _tb(s, Inches(1.0), Inches(2.0), Inches(9.0), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
_set(tf, "Resume Scrubber", sz=54, bold=True)
_p(tf, "Automated PII Redaction & CV Template Population", sz=22, col=DARK_GRAY, after=Pt(18))
_p(tf, "Quality & Compliance Engineering  ·  Gilead Sciences", sz=15, col=MID_GRAY)
_line(s, Inches(1.0), Inches(4.5), Inches(3.2))
_sn(s, 1)


# ═══════════════════════════════════════════════════════════════
# 2 — AGENDA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(4.0), Inches(0.7))
_set(tb.text_frame, "Agenda", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(1.6))

items = [
    ("01", "Problem Statement",     "Why this tool was needed"),
    ("02", "Initial Ideas",         "Copilot Studio vs. non-AI approach"),
    ("03", "Methods Tried",         "Evaluating trade-offs"),
    ("04", "Final Solution",        "Architecture & technical design"),
    ("05", "Cross-Functional Work", "R&D Training, admins, coordinators"),
    ("06", "HCI & Iteration",       "User feedback and UI evolution"),
    ("07", "Design Process",        "Stakeholder feedback, WalkMe alignment"),
    ("08", "Demo",                  "Live walkthrough of the tool"),
    ("09", "Lessons Learned",       "Experiences gained"),
]
for i, (num, title, desc) in enumerate(items):
    y = Inches(1.4 + i * 0.68)
    c = _oval(s, Inches(0.9), y + Pt(5), Inches(0.38), Inches(0.38),
              RED if i == 0 else LIGHT_GRAY)
    tb = _tb(s, Inches(0.9), y + Pt(5), Inches(0.38), Inches(0.38))
    _set(tb.text_frame, num, sz=11, bold=True, col=WHITE if i == 0 else DARK_GRAY,
         align=PP_ALIGN.CENTER)
    tb = _tb(s, Inches(1.55), y, Inches(5.5), Inches(0.55))
    tf = tb.text_frame
    _set(tf, title, sz=16, bold=True, after=Pt(1))
    _p(tf, desc, sz=11, col=MID_GRAY, after=Pt(0))
_sn(s, 2)


# ═══════════════════════════════════════════════════════════════
# 3 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.7))
_set(tb.text_frame, "Problem Statement", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(2.4))

pain = [
    ("PII Exposure Risk",      "Resumes contain emails, phone numbers, home\naddresses, photos, and document metadata"),
    ("Manual Redaction",        "Training coordinators spend hours per cycle\nreformatting and sanitising 50+ CVs"),
    ("Inconsistent Output",    "No standard process — each person redacts\ndifferently, creating compliance gaps"),
    ("Template Conformance",   "CVs must match FRM-11110 template format\nbefore upload to Veeva Vault (GVault)"),
]
for i, (t, d) in enumerate(pain):
    y = Inches(1.35 + i * 1.42)
    _rrect(s, Inches(0.7), y, Inches(5.9), Inches(1.2), LIGHT_GRAY)
    _oval(s, Inches(0.95), y + Inches(0.15), Inches(0.16), Inches(0.16), RED)
    tb = _tb(s, Inches(1.3), y + Inches(0.08), Inches(5.1), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, t, sz=15, bold=True, after=Pt(2))
    _p(tf, d, sz=12, col=DARK_GRAY, after=Pt(0))

# Impact metric cards (right)
metrics = [("50+", "CVs per cycle"), ("4–6 hrs", "Manual redaction\ntime per batch"), ("100%", "Must be PII-free\nfor GVault")]
for i, (big, lbl) in enumerate(metrics):
    y = Inches(1.35 + i * 1.88)
    _rrect(s, Inches(7.1), y, Inches(2.9), Inches(1.55), RED_LIGHT)
    tb = _tb(s, Inches(7.1), y + Inches(0.15), Inches(2.9), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, big, sz=38, bold=True, col=RED, align=PP_ALIGN.CENTER)
    _p(tf, lbl, sz=12, col=DARK_GRAY, align=PP_ALIGN.CENTER, after=Pt(0))
_sn(s, 3)


# ═══════════════════════════════════════════════════════════════
# 4 — INITIAL IDEAS: COMPARISON TABLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Initial Ideas — Copilot Studio vs. Non-AI", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(3.6))

col_w, x_lbl, x_a, x_b = Inches(3.5), Inches(0.8), Inches(4.6), Inches(8.3)
y0 = Inches(1.45)
_rrect(s, x_a, y0, col_w, Inches(0.58), RED)
tb = _tb(s, x_a, y0 + Pt(5), col_w, Inches(0.5))
_set(tb.text_frame, "Copilot Studio / LLM", sz=15, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
_rrect(s, x_b, y0, col_w, Inches(0.58), NEAR_BLACK)
tb = _tb(s, x_b, y0 + Pt(5), col_w, Inches(0.5))
_set(tb.text_frame, "Deterministic + Lightweight ML", sz=15, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Data Privacy",    "Resume sent to external LLM",       "All processing on-premises  ✓"),
    ("Output Control",  "Unpredictable OOXML formatting",    "Full XML control  ✓"),
    ("Latency",         "High — API round-trip per request",  "Low — local inference  ✓"),
    ("Cost",            "Per-call pricing at scale",          "No per-request cost  ✓"),
    ("PII Detection",   "Good recall, hallucination risk",   "Regex + validation gates  ✓"),
    ("Reproducibility", "Non-deterministic outputs",         "Deterministic  ✓"),
    ("Maintenance",     "Vendor-dependent updates",          "Self-managed  ✓"),
    ("Setup Effort",    "Quick prototype, low barrier",       "More initial engineering  △"),
]
for i, (lbl, a, b) in enumerate(rows):
    y = y0 + Inches(0.62 + i * 0.58)
    bg = WHITE if i % 2 == 0 else LIGHT_GRAY
    _rrect(s, x_lbl, y, Inches(3.5), Inches(0.52), bg)
    tb = _tb(s, x_lbl + Inches(0.15), y + Pt(4), Inches(3.2), Inches(0.45))
    _set(tb.text_frame, lbl, sz=12, bold=True)
    _rrect(s, x_a, y, col_w, Inches(0.52), bg)
    tb = _tb(s, x_a + Inches(0.1), y + Pt(4), Inches(3.3), Inches(0.45))
    _set(tb.text_frame, a, sz=11, col=DARK_GRAY, align=PP_ALIGN.CENTER)
    _rrect(s, x_b, y, col_w, Inches(0.52), bg)
    tb = _tb(s, x_b + Inches(0.1), y + Pt(4), Inches(3.3), Inches(0.45))
    _set(tb.text_frame, b, sz=11, col=DARK_GRAY, align=PP_ALIGN.CENTER)

# Selected badge
y_badge = y + Inches(0.7)
_rrect(s, x_b + Inches(0.5), y_badge, Inches(2.5), Inches(0.42), RED_LIGHT, line_c=RED)
tb = _tb(s, x_b + Inches(0.5), y_badge + Pt(2), Inches(2.5), Inches(0.38))
_set(tb.text_frame, "▶  Selected approach", sz=12, bold=True, col=RED, align=PP_ALIGN.CENTER)
_sn(s, 4)


# ═══════════════════════════════════════════════════════════════
# 5 — METHODS TRIED (card flow)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Methods Tried & Considered", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(3.0))

methods = [
    ("python-docx\n(High-level API)",   "Initial library used for all\ndocument manipulation.\nCan't inject XML across docs,\nmerge numbering, or manipulate\nraw OOXML (sdt, hyperlinks, rels)",    False),
    ("Header-Only\nRegex Parser",       "Dictionary of ~50 section aliases\nwith canonical normalization.\nFailed on resumes with no headers\nor creative/non-standard titles",      False),
    ("Pure spaCy\nNER Detection",       "Entity labels (DEGREE, COMPANY)\nclassify paragraphs.\nFinds entities, not boundaries —\nentity leakage across sections",    False),
    ("Cloud PII APIs\n(DLP/Comprehend)","Resume data leaves premises —\nunacceptable in regulated\n(GVault/Veeva) environment",               False),
    ("Hybrid ML +\nHeuristics + lxml",  "Headers → boundaries (precise)\nNER → entity votes (fills gaps)\nlxml → full OOXML control\nAll on-premises, zero API calls",         True),
]
cw = Inches(2.15); gap = Inches(0.25)
tw = len(methods) * cw + (len(methods)-1) * gap
sx = (SW - tw) // 2

for i, (name, reason, ok) in enumerate(methods):
    x = sx + i * (cw + gap); y = Inches(1.6)
    bg = GREEN_LIGHT if ok else RED_LIGHT
    _rrect(s, x, y, cw, Inches(3.6), bg, line_c=RED if ok else None)
    # badge
    bc = GREEN_DARK if ok else RED
    _oval(s, x + cw - Inches(0.52), y + Inches(0.08), Inches(0.38), Inches(0.38), bc)
    tb = _tb(s, x + cw - Inches(0.52), y + Inches(0.08), Inches(0.38), Inches(0.38))
    _set(tb.text_frame, "✓" if ok else "✗", sz=15, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    # name
    tb = _tb(s, x + Inches(0.12), y + Inches(0.55), cw - Inches(0.24), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, name, sz=13, bold=True, align=PP_ALIGN.CENTER, after=Pt(0))
    # reason
    tb = _tb(s, x + Inches(0.12), y + Inches(1.4), cw - Inches(0.24), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, reason, sz=10, col=DARK_GRAY, align=PP_ALIGN.CENTER, after=Pt(0))
    # arrow
    if i < len(methods)-1:
        _arrow(s, x + cw + Pt(2), y + Inches(1.6), gap - Pt(4), Inches(0.3), MID_GRAY)

# Insight bar
tb = _tb(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
_set(tf, "Key Insight", sz=16, bold=True, col=RED, after=Pt(4))
_p(tf, "Each approach revealed a critical limitation that shaped the next iteration. "
   "python-docx was the initial library but lacked raw XML control. "
   "Header-only regex was reliable but too rigid. Pure NER found entities without boundaries. "
   "Cloud APIs were a non-starter for data privacy. "
   "The final hybrid combines all strengths — on-premises with no external dependencies.", sz=13, col=DARK_GRAY, after=Pt(0))
_sn(s, 5)


# ═══════════════════════════════════════════════════════════════
# 5b — INITIAL REGEX APPROACH & RULES
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Initial Approach — Regex-Based Section Detection", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(4.0))

# Left column: Header-Only Rules
_rrect(s, Inches(0.7), Inches(1.4), Inches(5.8), Inches(5.6), LIGHT_GRAY)
_rect(s, Inches(0.7), Inches(1.4), Inches(5.8), Pt(3), fill=RED)
tb = _tb(s, Inches(0.9), Inches(1.55), Inches(5.4), Inches(0.4))
_set(tb.text_frame, "Header-Only Parser Rules", sz=15, bold=True, col=RED)

rules = [
    ("Section Keyword Matching", "Maintained dictionary of 50+ section aliases\n(\"Experience\", \"Work History\", \"Employment\", etc.)"),
    ("Canonical Normalization", "Collapsed spaced letters (E X P E R I E N C E),\nstripped non-alphanumerics, lowercased"),
    ("Uniqueness Gate", "Canonical text must appear exactly once in the\ndocument (prevents body text → header confusion)"),
    ("Length Guards", "Line ≤120 chars, canonical form ≤40 chars,\nmust not end with '.' or ':'"),
    ("Word Ratio Fallback", "≥50% of line words match a known keyword\nphrase — catches partial/fuzzy headers"),
]
for j, (rlabel, rdesc) in enumerate(rules):
    yy = Inches(2.0 + j * 0.98)
    _oval(s, Inches(0.95), yy + Inches(0.05), Inches(0.14), Inches(0.14), RED)
    tb = _tb(s, Inches(1.3), yy - Inches(0.03), Inches(5.0), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, rlabel, sz=12, bold=True, after=Pt(1))
    _p(tf, rdesc, sz=11, col=DARK_GRAY, after=Pt(0))

# Right column: Limitations
_rrect(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.6), RED_LIGHT)
_rect(s, Inches(6.8), Inches(1.4), Inches(5.8), Pt(3), fill=RED)
tb = _tb(s, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.4))
_set(tb.text_frame, "Why It Wasn't Enough", sz=15, bold=True, col=RED)

limits = [
    ("No headers at all", "Some resumes use formatting (bold, size)\ninstead of explicit section labels"),
    ("Non-standard headers", "Creative titles like \"Where I've Been\"\nor \"My Journey\" → not in any dictionary"),
    ("Paragraphs outside zones", "Content before the first header or between\nsections gets lost entirely"),
    ("No entity understanding", "Can't distinguish education from experience\nwhen headers are ambiguous or missing"),
    ("Verdict", "Reliable as a boundary signal but insufficient\nas the sole detection method → led to hybrid"),
]
for j, (rlabel, rdesc) in enumerate(limits):
    yy = Inches(2.0 + j * 0.98)
    _oval(s, Inches(7.05), yy + Inches(0.05), Inches(0.14), Inches(0.14), RED)
    tb = _tb(s, Inches(7.4), yy - Inches(0.03), Inches(5.0), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, rlabel, sz=12, bold=True, after=Pt(1))
    _p(tf, rdesc, sz=11, col=DARK_GRAY, after=Pt(0))
_sn(s, "5b")


# ═══════════════════════════════════════════════════════════════
# 5b2 — PII REGEX PATTERNS (detailed)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "PII Detection — Regex Patterns & Validation", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(4.0))

# Pattern cards in 2x3 grid
pii_cards = [
    ("Phone Number", RED, [
        "Lookbehind: (?<![\\w@.]) — skip emails/IDs",
        "Intl prefix: (+1, +44, 0044) optional",
        "Area code: (555), ( 201 ) optional",
        "Core: 2–6 digit groups with separators",
        "Extension: ext/x + digits optional",
        "Validation: 7–15 digits required (E.164)",
        "Rejects: \"2016\" (4 digits), \"Bachelor\"",
    ]),
    ("Email Address", RED, [
        "Pattern: user@domain.tld",
        "user: [a-zA-Z0-9._%+-]+",
        "domain: [a-zA-Z0-9.-]+",
        "TLD: [a-zA-Z]{2,}",
        "",
        "Standard RFC-compliant regex",
        "No validation gate needed",
    ]),
    ("URL / Website", RED, [
        "Matches: https://, http://, www.",
        "Also: bare domains with common TLDs",
        "TLDs: .com .org .net .io .dev .me",
        "        .co .info .biz",
        "Catches: linkedin.com/in/user",
        "         github.com/username",
        "No protocol prefix required",
    ]),
    ("City, STATE", NEAR_BLACK, [
        "Pattern: City Name, XX",
        "City: capital letter + 1–30 chars",
        "State: US (50 + DC), AU (8), CA (13)",
        "Examples:",
        "  San Francisco, CA  ✓",
        "  New York, NY  ✓",
        "  experience, management  ✗",
    ]),
    ("Address (Multi-Signal)", NEAR_BLACK, [
        "Postal codes: US/UK/CA/NL/JP (+3 pts)",
        "Street types: 10+ languages (+3 pts)",
        "House number: 1–5 digits (+2 pts)",
        "Region codes: state abbrevs (+2 pts)",
        "Unit/suite, direction (+1 pt each)",
        "PO Box, country name (+1–3 pts)",
        "Threshold: score ≥ 5 → address",
    ]),
    ("PII-Safe Address Check", NEAR_BLACK, [
        "Problem: phone digits inflate score",
        "  (201) 320-6355 → 320-6355 matches",
        "  Japanese postal code pattern \\d{3}-\\d{4}",
        "",
        "Fix: strip phone/email/URL from text",
        "BEFORE running address scorer",
        "Then score remaining text only",
    ]),
]

pcw = Inches(3.85); pg = Inches(0.15)
for i, (ptitle, accent, pitems) in enumerate(pii_cards):
    col = i % 3
    row = i // 3
    x = Inches(0.7) + col * (pcw + pg)
    y = Inches(1.4) + row * Inches(3.05)
    _rrect(s, x, y, pcw, Inches(2.85), LIGHT_GRAY)
    _rrect(s, x, y, pcw, Inches(0.42), accent)
    tb = _tb(s, x, y + Inches(0.05), pcw, Inches(0.32))
    _set(tb.text_frame, ptitle, sz=13, bold=True, col=WHITE, align=PP_ALIGN.CENTER, after=Pt(0))
    for j, pitem in enumerate(pitems):
        tb = _tb(s, x + Inches(0.12), y + Inches(0.48 + j * 0.32), pcw - Inches(0.24), Inches(0.3))
        tf = tb.text_frame; tf.word_wrap = True
        _set(tf, pitem, sz=9, col=DARK_GRAY if pitem else WHITE, after=Pt(0))
_sn(s, "5b2")


# ═══════════════════════════════════════════════════════════════
# 5c — spaCy NER MODEL
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "NER Model — spaCy & Fine-Tuning", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(3.2))

# Three cards across
cards = [
    ("Pretrained Base Model", RED, [
        ("Model", "en_core_web_md (spaCy 3.8)"),
        ("Architecture", "tok2vec + tagger + parser +\nsenter + NER + attribute_ruler"),
        ("Vectors", "300-dimensional word vectors\n(20,000 unique vectors, 684K keys)"),
        ("License", "MIT (Explosion AI)"),
        ("Why this model?", "Medium-sized — good balance of\naccuracy and inference speed for\non-premises CPU deployment"),
    ]),
    ("Fine-Tuning Dataset", NEAR_BLACK, [
        ("Source", "Resume Entities for NER\n(Kaggle public dataset)"),
        ("Format", "Annotated resume text with\nentity spans and labels"),
        ("Entity Labels", "COLLEGE_NAME, DEGREE,\nGRADUATION_YEAR,\nCOMPANIES_WORKED_AT,\nYEARS_OF_EXPERIENCE"),
        ("Training", "spaCy CLI training pipeline\nwith config.cfg"),
        ("Purpose", "Teach model to recognise resume-\nspecific entities (not general NER)"),
    ]),
    ("Hybrid Integration", DARK_GRAY, [
        ("Role", "Secondary signal — fills gaps\nwhere headers are missing"),
        ("Voting", "EDU entities → \"education\" vote\nEXP entities → \"experience\" vote"),
        ("Resolution", "Headers always win when present;\nmodel fills header-less zones"),
        ("Block Propagation", "Unvoted paragraphs inherit the\nsection of surrounding voted ones"),
        ("Fallback", "If model fails to load, degrades\ngracefully to header-only parser"),
    ]),
]
ccw = Inches(3.85); cg = Inches(0.15)
ctw = len(cards)*ccw + (len(cards)-1)*cg
csx = (SW - ctw) // 2

for ci, (heading, accent, items) in enumerate(cards):
    x = csx + ci*(ccw+cg); y = Inches(1.45)
    _rrect(s, x, y, ccw, Inches(5.6), LIGHT_GRAY)
    _rrect(s, x, y, ccw, Inches(0.5), accent)
    tb = _tb(s, x, y + Inches(0.06), ccw, Inches(0.38))
    _set(tb.text_frame, heading, sz=14, bold=True, col=WHITE, align=PP_ALIGN.CENTER, after=Pt(0))
    for j, (lbl, val) in enumerate(items):
        iy = y + Inches(0.6 + j * 0.95)
        tb = _tb(s, x + Inches(0.15), iy, ccw - Inches(0.3), Inches(0.9))
        tf = tb.text_frame; tf.word_wrap = True
        _set(tf, lbl, sz=11, bold=True, after=Pt(1))
        _p(tf, val, sz=10, col=DARK_GRAY, after=Pt(0))
_sn(s, "5c")


# ═══════════════════════════════════════════════════════════════
# 5d — CURRENT TECH STACK
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Current Technology Stack", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(3.0))

stack_rows = [
    ("Frontend",        "React 18, Vite 5, React Quill 2.0",     "Step-by-step wizard UI, rich text editor\nfor responsibilities, drag-and-drop upload"),
    ("API Layer",       "Flask, Flask-CORS, Gunicorn",            "Two endpoints: /remove-images (PII redaction),\n/populate-template (CV template). 120s timeout"),
    ("Document I/O",    "python-docx",                            "Open/save .docx, create paragraphs,\nread resolved font properties (name, size)"),
    ("XML Engine",      "lxml (etree)",                           "Direct OOXML surgery: PII span redaction,\nhyperlink unwrapping, cross-doc XML injection"),
    ("NLP / ML",        "spaCy 3.8 (en_core_web_md fine-tuned)", "NER entity voting for section classification.\nCPU inference, ~200ms per resume"),
    ("PII Detection",   "Regex + validation gates",              "Phone, email, URL, city/state, address scoring.\nNo external API calls — all on-premises"),
    ("HTML → OOXML",    "html_to_docx.py (parse_quill_html)",    "Converts Quill rich text HTML to OOXML\nparagraph specs (bold, italic, underline, lists)"),
    ("Deployment",      "Google App Engine (app.yaml)",           "On-premises hosting, zero external\nAPI dependencies, single-instance"),
]

# Table header
_rrect(s, Inches(0.7), Inches(1.4), Inches(2.8), Inches(0.5), RED)
tb = _tb(s, Inches(0.7), Inches(1.43), Inches(2.8), Inches(0.44))
_set(tb.text_frame, "Layer", sz=13, bold=True, col=WHITE, align=PP_ALIGN.CENTER, after=Pt(0))
_rrect(s, Inches(3.5), Inches(1.4), Inches(3.8), Inches(0.5), RED)
tb = _tb(s, Inches(3.5), Inches(1.43), Inches(3.8), Inches(0.44))
_set(tb.text_frame, "Technology", sz=13, bold=True, col=WHITE, align=PP_ALIGN.CENTER, after=Pt(0))
_rrect(s, Inches(7.3), Inches(1.4), Inches(5.3), Inches(0.5), RED)
tb = _tb(s, Inches(7.3), Inches(1.43), Inches(5.3), Inches(0.44))
_set(tb.text_frame, "Role", sz=13, bold=True, col=WHITE, align=PP_ALIGN.CENTER, after=Pt(0))

for i, (layer, tech, role) in enumerate(stack_rows):
    y = Inches(1.9 + i * 0.68)
    bg = WHITE if i % 2 == 0 else LIGHT_GRAY
    _rrect(s, Inches(0.7), y, Inches(2.8), Inches(0.62), bg)
    tb = _tb(s, Inches(0.85), y + Pt(4), Inches(2.5), Inches(0.55))
    _set(tb.text_frame, layer, sz=11, bold=True, after=Pt(0))
    _rrect(s, Inches(3.5), y, Inches(3.8), Inches(0.62), bg)
    tb = _tb(s, Inches(3.6), y + Pt(4), Inches(3.6), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, tech, sz=10, col=DARK_GRAY, after=Pt(0))
    _rrect(s, Inches(7.3), y, Inches(5.3), Inches(0.62), bg)
    tb = _tb(s, Inches(7.4), y + Pt(4), Inches(5.1), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, role, sz=10, col=DARK_GRAY, after=Pt(0))
_sn(s, "5d")


# ═══════════════════════════════════════════════════════════════
# 6 — FINAL SOLUTION: ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Final Solution — Architecture", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(3.0))

# Frontend layer
_rrect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.85), RED_LIGHT)
tb = _tb(s, Inches(1.1), Inches(1.55), Inches(2.0), Inches(0.4))
_set(tb.text_frame, "Frontend", sz=14, bold=True, col=RED)
tb = _tb(s, Inches(3.0), Inches(1.58), Inches(9.0), Inches(0.35))
_set(tb.text_frame, "React 18  ·  Vite 5  ·  React Quill  ·  Drag-and-drop upload  ·  Step-by-step wizard", sz=12, col=DARK_GRAY)

# API layer
_rrect(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.85), LIGHT_GRAY)
tb = _tb(s, Inches(1.1), Inches(2.65), Inches(2.0), Inches(0.4))
_set(tb.text_frame, "API Layer", sz=14, bold=True, col=NEAR_BLACK)
tb = _tb(s, Inches(3.0), Inches(2.68), Inches(9.0), Inches(0.35))
_set(tb.text_frame, "Flask  ·  Gunicorn (120s timeout)  ·  Flask-CORS  ·  Quill HTML → OOXML converter", sz=12, col=DARK_GRAY)

# Arrows down
_arrow_down_pts = [(Inches(3.2), "POST /remove-images"), (Inches(9.2), "POST /populate-template")]
for ax, lbl in _arrow_down_pts:
    _rect(s, ax, Inches(3.45), Pt(3), Inches(0.35), fill=RED)
    # small triangle
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, ax - Inches(0.08), Inches(3.75), Inches(0.2), Inches(0.15))
    tri.fill.solid(); tri.fill.fore_color.rgb = RED; tri.line.fill.background()
    tri.rotation = 180.0
    tb = _tb(s, ax - Inches(1.2), Inches(3.5), Inches(2.4), Inches(0.3))
    _set(tb.text_frame, lbl, sz=10, bold=True, col=RED, align=PP_ALIGN.CENTER)

# Workflow A
wf_items = [
    ("Workflow A — PII Redaction", RED, [
        ("clean_resume.py", "Regex PII + City/STATE + URL detection"),
        ("address_identifier.py", "Multilingual address scoring"),
        ("_insert_experience_entry", "Gilead role insertion (3-step fallback)"),
        ("Image & hyperlink removal", "OOXML tree surgery + metadata scrub"),
    ]),
    ("Workflow B — Template Population", NEAR_BLACK, [
        ("parser_get_text.py", "Extract (text, XML) pairs"),
        ("model_section_parser.py", "Hybrid ML + header detection"),
        ("parser_get_section_xml.py", "Sanitise & deep-copy XML"),
        ("populate_template.py", "OOXML injection + invariants"),
    ]),
]
wx = [Inches(0.8), Inches(6.55)]
for idx, (title, accent, items) in enumerate(wf_items):
    x = wx[idx]; y = Inches(4.1)
    _rrect(s, x, y, Inches(5.5), Inches(3.0), LIGHT_GRAY)
    _rect(s, x, y, Inches(5.5), Pt(4), fill=accent)
    tb = _tb(s, x + Inches(0.2), y + Inches(0.15), Inches(5.0), Inches(0.35))
    _set(tb.text_frame, title, sz=13, bold=True, col=accent)
    for j, (mod, desc) in enumerate(items):
        yy = y + Inches(0.55 + j * 0.55)
        tb = _tb(s, x + Inches(0.25), yy, Inches(2.3), Inches(0.35))
        _set(tb.text_frame, mod, sz=11, bold=True, col=NEAR_BLACK, after=Pt(0))
        tb = _tb(s, x + Inches(2.6), yy, Inches(2.7), Inches(0.35))
        _set(tb.text_frame, desc, sz=11, col=DARK_GRAY, after=Pt(0))
_sn(s, 6)


# ═══════════════════════════════════════════════════════════════
# 7 — DATA FLOW: CHEVRON PIPELINE + DETAIL CARDS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Template Population — Data Flow", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(3.2))

stages = [
    ("1", "Extract",  "Text + XML\npairs from .docx"),
    ("2", "Detect",   "Hybrid section\nclassification"),
    ("3", "Build",    "Deep-copy &\nsanitise XML"),
    ("4", "Inject",   "Merge into\ntemplate"),
    ("5", "Replace",  "User info +\nQuill → OOXML"),
    ("6", "Deliver",  "Return .docx\ndownload"),
]
cw2 = Inches(1.75); gap2 = Inches(0.18)
tw2 = len(stages)*cw2 + (len(stages)-1)*gap2
sx2 = (SW - tw2) // 2
yc = Inches(1.5)

for i, (n, lbl, desc) in enumerate(stages):
    x = sx2 + i*(cw2+gap2)
    _chevron(s, x, yc, cw2, Inches(1.15), RED if i == 0 else LIGHT_GRAY)
    tb = _tb(s, x + Inches(0.28), yc + Inches(0.15), cw2 - Inches(0.4), Inches(0.35))
    _set(tb.text_frame, f"{n}. {lbl}", sz=12, bold=True,
         col=WHITE if i == 0 else NEAR_BLACK, align=PP_ALIGN.CENTER, after=Pt(0))
    tb = _tb(s, x + Inches(0.05), yc + Inches(1.25), cw2 - Inches(0.1), Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, desc, sz=10, col=DARK_GRAY, align=PP_ALIGN.CENTER, after=Pt(0))

# Divider
_rect(s, Inches(0.8), Inches(3.9), Inches(11.7), Pt(1), fill=LIGHT_GRAY)

# Detail cards
det = [
    ("Section Detection (Hybrid)", [
        "Pass 1 — Header keywords → hard zone boundaries",
        "Pass 2 — spaCy NER → entity votes per paragraph",
        "Resolution — Headers win; model fills gaps",
        "Fallback — Header-only if model fails to load",
    ]),
    ("PII Redaction (Multi-Stage)", [
        "Phone: regex → digit-count validation (7–15)",
        "Email: standard regex pattern matching",
        "URL: bare domains (linkedin.com) + protocols",
        "City, STATE: pattern match (San Francisco, CA)",
        "Address: postal + street + house # scoring",
    ]),
    ("OOXML Invariant Enforcement", [
        "No empty <w:sdtContent/> — Word rejects these",
        "Every <w:tc> must contain ≥1 <w:p>",
        "Numbering defs merged & deduplicated",
        "Style refs stripped to avoid ID mismatches",
    ]),
]
dcw = Inches(3.7); dg = Inches(0.2)
dtw = len(det)*dcw + (len(det)-1)*dg
dsx = (SW - dtw) // 2

for i, (title, items) in enumerate(det):
    x = dsx + i*(dcw+dg); y = Inches(4.1)
    _rrect(s, x, y, dcw, Inches(3.0), LIGHT_GRAY)
    _rect(s, x, y, dcw, Pt(3), fill=RED)
    tb = _tb(s, x + Inches(0.2), y + Inches(0.15), dcw - Inches(0.4), Inches(0.35))
    _set(tb.text_frame, title, sz=13, bold=True, col=RED, after=Pt(0))
    for j, item in enumerate(items):
        tb = _tb(s, x + Inches(0.2), y + Inches(0.55 + j*0.55), dcw - Inches(0.4), Inches(0.5))
        tf = tb.text_frame; tf.word_wrap = True
        _set(tf, "·  " + item, sz=11, col=DARK_GRAY, after=Pt(0))
_sn(s, 7)


# ═══════════════════════════════════════════════════════════════
# 8 — CROSS-FUNCTIONAL COLLABORATION (hub-spoke)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Cross-Functional Collaboration", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(3.2))

# Hub
_oval(s, Inches(5.4), Inches(2.8), Inches(2.5), Inches(2.5), RED)
htb = _tb(s, Inches(5.4), Inches(3.45), Inches(2.5), Inches(1.2))
htf = htb.text_frame; htf.word_wrap = True
_set(htf, "Resume\nScrubber", sz=18, bold=True, col=WHITE, align=PP_ALIGN.CENTER, after=Pt(0))

spokes = [
    ("R&D Training Team", ["Understood CV submission workflow",
                           "Validated FRM-11110 template",
                           "Provided diverse sample resumes"], Inches(0.6), Inches(1.4)),
    ("Training Admins\n& Coordinators", ["Revealed 50+ CVs per cycle",
                                          "Identified metadata scrubbing gaps",
                                          "Shaped the ≥150 char requirement"], Inches(0.6), Inches(4.5)),
    ("Quality &\nCompliance",  ["Confirmed PII redaction requirements",
                                 "Validated regex detection thresholds",
                                 "Approved post-processing warning"], Inches(8.9), Inches(1.4)),
    ("End Users",              ["Tested iterative UI prototypes",
                                 "Surfaced real-world edge cases",
                                 "Became early adopters & evangelists"], Inches(8.9), Inches(4.5)),
]
for title, items, x, y in spokes:
    _rrect(s, x, y, Inches(3.9), Inches(2.2), LIGHT_GRAY)
    _rect(s, x, y, Inches(3.9), Pt(3), fill=RED)
    tb = _tb(s, x + Inches(0.2), y + Inches(0.12), Inches(3.5), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, title, sz=13, bold=True, after=Pt(2))
    for j, item in enumerate(items):
        tb = _tb(s, x + Inches(0.2), y + Inches(0.7 + j*0.42), Inches(3.5), Inches(0.38))
        tf = tb.text_frame; tf.word_wrap = True
        _set(tf, "·  " + item, sz=11, col=DARK_GRAY, after=Pt(0))

# Connectors (simple lines)
_rect(s, Inches(4.5), Inches(2.4), Inches(1.1), Pt(2), fill=MID_GRAY)
_rect(s, Inches(4.5), Inches(5.3), Inches(1.1), Pt(2), fill=MID_GRAY)
_rect(s, Inches(7.7), Inches(2.4), Inches(1.2), Pt(2), fill=MID_GRAY)
_rect(s, Inches(7.7), Inches(5.3), Inches(1.2), Pt(2), fill=MID_GRAY)
_sn(s, 8)


# ═══════════════════════════════════════════════════════════════
# 9 — HCI & ITERATIVE DESIGN (timeline)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "HCI & Iterative Design", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(2.8))

iterations = [
    ("Iteration 1", "CLI Prototype",
     ["Command-line tool for testing", "Proved pipeline end-to-end",
      "Feedback: non-technical\nusers can't use this"]),
    ("Iteration 2", "Basic Web Form",
     ["Simple upload → download", "Missing user info fields",
      "Feedback: 'Where do I\nenter my info?'"]),
    ("Iteration 3", "Guided Steps",
     ["Card layout with 5 steps", "Rich text editor added",
      "Drag-and-drop with\nvisual feedback"]),
    ("Iteration 4", "Production UI",
     ["Post-download PII warning", "Inline validation + error scroll",
      "Char count, GVault link,\nform polish"]),
]

# Timeline line
_rect(s, Inches(1.0), Inches(2.2), Inches(11.3), Pt(3), fill=RED)

cw3 = Inches(2.65); gap3 = Inches(0.28)
tw3 = len(iterations)*cw3 + (len(iterations)-1)*gap3
sx3 = (SW - tw3) // 2

for i, (phase, title, items) in enumerate(iterations):
    x = sx3 + i*(cw3+gap3)
    is_last = (i == len(iterations)-1)
    # dot
    dsz = Inches(0.26)
    dot = _oval(s, x + cw3/2 - dsz/2, Inches(2.08), dsz, dsz,
                RED if is_last else WHITE)
    if not is_last:
        dot.line.color.rgb = RED; dot.line.width = Pt(2)
    # phase label
    tb = _tb(s, x, Inches(1.6), cw3, Inches(0.4))
    _set(tb.text_frame, phase, sz=11, bold=True, col=RED, align=PP_ALIGN.CENTER, after=Pt(0))
    # card
    bg = RED_LIGHT if is_last else LIGHT_GRAY
    _rrect(s, x, Inches(2.6), cw3, Inches(3.3), bg, line_c=RED if is_last else None)
    tb = _tb(s, x + Inches(0.1), Inches(2.72), cw3 - Inches(0.2), Inches(0.4))
    _set(tb.text_frame, title, sz=14, bold=True,
         col=RED if is_last else NEAR_BLACK, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        tb = _tb(s, x + Inches(0.1), Inches(3.2 + j*0.7), cw3 - Inches(0.2), Inches(0.65))
        tf = tb.text_frame; tf.word_wrap = True
        _set(tf, item, sz=11, col=DARK_GRAY, align=PP_ALIGN.CENTER, after=Pt(0))

# Feedback annotation
tb = _tb(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
_set(tf, "Each iteration was shaped by direct feedback from training coordinators, admins, and end users. "
     "Continuous user testing drove every design decision — from the CLI prototype to the guided step-by-step production UI.",
     sz=13, col=DARK_GRAY, align=PP_ALIGN.CENTER, after=Pt(0))
_sn(s, 9)


# ═══════════════════════════════════════════════════════════════
# 9b — DESIGN PROCESS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Design Process — Stakeholder-Driven Development", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(4.2))

# Timeline line
_rect(s, Inches(1.0), Inches(2.1), Inches(11.3), Pt(3), fill=RED)

# Phase dots and cards
phases = [
    ("Phase 1", "Internal Iteration",
     "with Lizzy", [
         "Iterative prototyping sessions",
         "Tested with diverse resume formats",
         "Refined PII detection thresholds",
         "Validated experience section insertion",
     ]),
    ("Phase 2", "R&D Training Team",
     "Demo & Feedback", [
         "Demonstrated tool to R&D Training",
         "Validated FRM-11110 template workflow",
         "Feedback shaped field requirements\n(responsibilities, start date)",
         "Ensured alignment with CV submission\nprocess for easy adoption",
     ]),
    ("Phase 3", "Quality Team",
     "Review & Approval", [
         "Presented PII redaction capabilities",
         "Verified compliance requirements met",
         "Confirmed post-processing warning\nfor manual review step",
         "Approved for production use in\nregulated environment",
     ]),
    ("Phase 4", "WalkMe Integration",
     "with Philip & Darris", [
         "Collaborated on guided walkthrough",
         "Aligned WalkMe flow with existing\nWalkMe processes and standards",
         "Ensured step-by-step guidance\nmatches the tool's UI wizard",
         "Seamless onboarding experience\nfor new users",
     ]),
]

pcw = Inches(2.75); pg = Inches(0.2)
ptw = len(phases)*pcw + (len(phases)-1)*pg
psx = (SW - ptw) // 2

for i, (phase, title, subtitle, items) in enumerate(phases):
    x = psx + i*(pcw+pg)
    is_last = (i == len(phases)-1)
    # dot
    dsz = Inches(0.26)
    dot = _oval(s, x + pcw/2 - dsz/2, Inches(1.97), dsz, dsz,
                RED if is_last else WHITE)
    if not is_last:
        dot.line.color.rgb = RED; dot.line.width = Pt(2)
    # phase label
    tb = _tb(s, x, Inches(1.5), pcw, Inches(0.4))
    _set(tb.text_frame, phase, sz=11, bold=True, col=RED, align=PP_ALIGN.CENTER, after=Pt(0))
    # card
    bg = RED_LIGHT if is_last else LIGHT_GRAY
    _rrect(s, x, Inches(2.5), pcw, Inches(4.3), bg, line_c=RED if is_last else None)
    tb = _tb(s, x + Inches(0.1), Inches(2.6), pcw - Inches(0.2), Inches(0.35))
    _set(tb.text_frame, title, sz=14, bold=True,
         col=RED if is_last else NEAR_BLACK, align=PP_ALIGN.CENTER, after=Pt(0))
    tb = _tb(s, x + Inches(0.1), Inches(2.95), pcw - Inches(0.2), Inches(0.3))
    _set(tb.text_frame, subtitle, sz=11, col=MID_GRAY, align=PP_ALIGN.CENTER, after=Pt(0))
    for j, item in enumerate(items):
        tb = _tb(s, x + Inches(0.1), Inches(3.4 + j*0.8), pcw - Inches(0.2), Inches(0.75))
        tf = tb.text_frame; tf.word_wrap = True
        _set(tf, "·  " + item, sz=11, col=DARK_GRAY, align=PP_ALIGN.LEFT, after=Pt(0))

# Bottom insight
tb = _tb(s, Inches(1.0), Inches(7.0), Inches(11.3), Inches(0.4))
tf = tb.text_frame; tf.word_wrap = True
_set(tf, "Every phase ensured the tool was tailored to existing team processes — enabling easy adoption without workflow disruption.",
     sz=13, col=DARK_GRAY, align=PP_ALIGN.CENTER, after=Pt(0))
_sn(s, "9b")


# ═══════════════════════════════════════════════════════════════
# 10 — LESSONS LEARNED (three columns)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Experiences Gained & Lessons Learned", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(3.8))

cols = [
    ("Technical", RED, [
        "OOXML direct manipulation gives\ncontrol but demands schema knowledge",
        "Hybrid ML + heuristics outperforms\neither approach alone",
        "Regex PII detection with validation\ngates balances precision & recall",
        "Pair-based (text, XML) architecture\nprevents alignment bugs",
        "Always build graceful degradation\npaths — have a fallback",
    ]),
    ("Cross-Functional", NEAR_BLACK, [
        "Talk to users early and often —\nassumptions about workflows are wrong",
        "Non-technical stakeholders provide\nthe most valuable feature insights",
        "Domain expertise from training teams\nimproved detection accuracy",
        "Early stakeholder alignment prevents\nrework late in development",
        "Coordinators became early adopters\nand evangelists for the tool",
    ]),
    ("Professional Growth", DARK_GRAY, [
        "Strengthened skills in NLP, document\nprocessing, and web application design",
        "Gained cross-team collaboration\nexperience across departments",
        "Learned to balance technical ambition\nwith practical user needs",
        "Iterative design with continuous\nfeedback converges faster",
        "The best tool is the one people\nactually want to use — UX matters",
    ]),
]
lcw = Inches(3.7); lg = Inches(0.2)
ltw = len(cols)*lcw + (len(cols)-1)*lg
lsx = (SW - ltw) // 2

for ci, (heading, accent, items) in enumerate(cols):
    x = lsx + ci*(lcw+lg); y = Inches(1.4)
    _rrect(s, x, y, lcw, Inches(0.5), accent)
    tb = _tb(s, x, y + Inches(0.06), lcw, Inches(0.38))
    _set(tb.text_frame, heading, sz=14, bold=True, col=WHITE, align=PP_ALIGN.CENTER, after=Pt(0))
    for j, item in enumerate(items):
        iy = y + Inches(0.65 + j*1.05)
        _rrect(s, x, iy, lcw, Inches(0.9), LIGHT_GRAY)
        _rect(s, x, iy, Pt(4), Inches(0.9), fill=accent)
        tb = _tb(s, x + Inches(0.2), iy + Inches(0.06), lcw - Inches(0.3), Inches(0.78))
        tf = tb.text_frame; tf.word_wrap = True
        _set(tf, item, sz=11, col=DARK_GRAY, after=Pt(0))
_sn(s, 10)


# ═══════════════════════════════════════════════════════════════
# 11 — IMPACT SUMMARY (with chart)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s); _sidebar(s)
tb = _tb(s, Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.7))
_set(tb.text_frame, "Impact Summary", sz=34, bold=True)
_line(s, Inches(0.8), Inches(1.0), Inches(2.0))

# Bar chart
cd = CategoryChartData()
cd.categories = ['Time per CV\n(minutes)', 'PII Detection\nCoverage (%)', 'Template\nConformance (%)']
cd.add_series('Before (Manual)', (25, 70, 50))
cd.add_series('After (Resume Scrubber)', (2, 95, 100))

cf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.4),
                         Inches(7.0), Inches(5.3), cd)
ch = cf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11); ch.legend.font.color.rgb = DARK_GRAY

s0 = ch.series[0]; s1 = ch.series[1]
s0.format.fill.solid(); s0.format.fill.fore_color.rgb = MID_GRAY
s1.format.fill.solid(); s1.format.fill.fore_color.rgb = RED

for sr, c in [(s0, DARK_GRAY), (s1, RED)]:
    sr.has_data_labels = True
    sr.data_labels.font.size = Pt(11); sr.data_labels.font.bold = True
    sr.data_labels.font.color.rgb = c
    sr.data_labels.number_format = '0'
    sr.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

ch.category_axis.tick_labels.font.size = Pt(12)
ch.category_axis.tick_labels.font.color.rgb = DARK_GRAY
ch.category_axis.has_major_gridlines = False
ch.value_axis.visible = False; ch.value_axis.has_major_gridlines = False
ch.plots[0].gap_width = 120

# Metric cards right
met = [
    ("12×", "faster",      "CV processing time\n~25 min → ~2 min"),
    ("95%", "coverage",    "Automated PII detection\nacross all categories"),
    ("100%", "conformance","Template output GVault-\nready without reformatting"),
    ("0",   "external calls","All data on-premises;\nno cloud API dependencies"),
]
for i, (big, lbl, desc) in enumerate(met):
    y = Inches(1.4 + i*1.4)
    _rrect(s, Inches(8.3), y, Inches(4.5), Inches(1.15), LIGHT_GRAY)
    _rect(s, Inches(8.3), y, Pt(4), Inches(1.15), fill=RED)
    tb = _tb(s, Inches(8.55), y + Inches(0.08), Inches(1.3), Inches(0.9))
    tf = tb.text_frame
    _set(tf, big, sz=30, bold=True, col=RED, after=Pt(0))
    _p(tf, lbl, sz=10, col=MID_GRAY, after=Pt(0))
    tb = _tb(s, Inches(10.0), y + Inches(0.15), Inches(2.6), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, desc, sz=11, col=DARK_GRAY, after=Pt(0))
_sn(s, 11)


# ═══════════════════════════════════════════════════════════════
# 11b — DEMO
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
_sidebar(s, Inches(0.18))
_oval(s, Inches(10.2), Inches(4.5), Inches(2.8), Inches(2.8), RED_LIGHT)
_oval(s, Inches(10.8), Inches(5.1), Inches(1.6), Inches(1.6), RED)

tb = _tb(s, Inches(1.0), Inches(2.0), Inches(9.0), Inches(2.0))
tf = tb.text_frame; tf.word_wrap = True
_set(tf, "Live Demo", sz=52, bold=True)
_p(tf, "Resume Scrubber in action", sz=22, col=DARK_GRAY, after=Pt(24))
_p(tf, "·  Upload a resume  ·  Enter Gilead role details  ·  Choose output format  ·  Download", sz=15, col=MID_GRAY, after=Pt(8))
_line(s, Inches(1.0), Inches(4.5), Inches(3.2))

# Two workflow cards
for i, (wf, desc) in enumerate([
    ("Keep My Resume Format", "PII redaction + Gilead experience\nentry insertion at section start"),
    ("Use CV Template", "Section extraction via hybrid ML +\ntemplate population (FRM-11110)"),
]):
    x = Inches(1.0 + i * 4.6); y = Inches(5.2)
    _rrect(s, x, y, Inches(4.2), Inches(1.6), LIGHT_GRAY)
    _rect(s, x, y, Inches(4.2), Pt(3), fill=RED)
    tb = _tb(s, x + Inches(0.2), y + Inches(0.15), Inches(3.8), Inches(0.35))
    _set(tb.text_frame, wf, sz=14, bold=True, col=RED)
    tb = _tb(s, x + Inches(0.2), y + Inches(0.55), Inches(3.8), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    _set(tf, desc, sz=12, col=DARK_GRAY, after=Pt(0))
_sn(s, "demo")


# ═══════════════════════════════════════════════════════════════
# 12 — THANK YOU
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
_oval(s, Inches(-1.5), Inches(-1.5), Inches(5.5), Inches(5.5), RED_LIGHT)
_oval(s, Inches(-0.5), Inches(-0.5), Inches(3.5), Inches(3.5), RED)
tb = _tb(s, Inches(4.5), Inches(2.4), Inches(7.0), Inches(1.8))
tf = tb.text_frame
_set(tf, "Thank You", sz=52, bold=True)
_p(tf, "Questions?", sz=24, col=MID_GRAY, after=Pt(20))
_p(tf, "Resume Scrubber  ·  Quality & Compliance Engineering", sz=14, col=MID_GRAY)
_line(s, Inches(4.5), Inches(2.25), Inches(2.5))
_sn(s, 12)


# ── Save ───────────────────────────────────────────────────────
out = "Resume_Scrubber_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
